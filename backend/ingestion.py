# WHAT DOES THIS FILE DO: Validate, extract, and chunk text from uploaded files for RAG ingestion

# ================== IMPORTS ==================
import csv
import io
import json
import os
from typing import Dict, List, Optional, Tuple

try:
    import fitz
    HAS_PDF = True
except Exception:
    HAS_PDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except Exception:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

try:
    from PIL import Image
    import pytesseract
    HAS_OCR = True
except Exception:
    HAS_OCR = False
# ================== IMPORTS ==================


# =========== VARIABLES : file configuration ===========
ALLOWED_EXTENSIONS = {
    ".txt", ".md", ".json", ".csv",
    ".pdf", ".doc", ".docx",
    ".ppt", ".pptx",
    ".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp",
}

_SIDECAR_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads', 'extracted_texts')
# =========== VARIABLES : file configuration ===========


# =========== FUNCTION ===========
# ROLE: Extract file extension from filename
def extension_of(filename: str) -> str:
    ''' Return lowercase file extension or empty string '''

    # FLOW-1: Return empty string if filename missing
    if not filename:
        return ""

    # FLOW-2: Split filename and return extension
    _, ext = os.path.splitext(filename.lower())
    return ext
# =========== FUNCTION ===========


# =========== FUNCTION ===========
# ROLE: Validate uploaded file type and size
def validate_upload(filename: str, file_size_bytes: int, max_size_bytes: int) -> None:
    ''' Raise ValueError if file is invalid, otherwise return silently '''

    # FLOW-1: Check file extension is supported
    ext = extension_of(filename)
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext or 'unknown'}")

    # FLOW-2: Check file is not empty
    if file_size_bytes <= 0:
        raise ValueError("Uploaded file is empty")

    # FLOW-3: Check file does not exceed size limit
    if file_size_bytes > max_size_bytes:
        raise ValueError(f"File too large: {file_size_bytes} bytes (max {max_size_bytes})")
# =========== FUNCTION ===========


# =========== FUNCTION ===========
# ROLE: Extract text from uploaded file bytes based on file type
def extract_text_from_bytes(filename: str, data: bytes) -> Tuple[str, Dict[str, str]]:
    ''' Return extracted text and metadata dict, raise ValueError if unsupported '''

    # FLOW-1: Determine file type and initialize metadata
    ext = extension_of(filename)
    metadata = {"extension": ext}

    # FLOW-2: Handle plain text and markdown files
    if ext in {".txt", ".md"}:
        return data.decode("utf-8", errors="ignore"), metadata

    # FLOW-3: Handle JSON files
    if ext == ".json":
        parsed = json.loads(data.decode("utf-8", errors="ignore"))
        return json.dumps(parsed, ensure_ascii=False, indent=2), metadata

    # FLOW-4: Handle CSV files — convert rows to pipe-separated format
    if ext == ".csv":
        text = data.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            rows.append(" | ".join(cell.strip() for cell in row if cell is not None))

        return "\n".join(r for r in rows if r.strip()), metadata

    # FLOW-5: Handle PDF files — extract text per page
    if ext == ".pdf":
        if not HAS_PDF:
            raise ValueError("PDF support unavailable. Install PyMuPDF.")
        doc = fitz.open(stream=data, filetype="pdf")
        pages = []
        for idx, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text and text.strip():
                pages.append(f"[Page {idx}]\n{text.strip()}")
        doc.close()

        return "\n\n".join(pages), metadata

    # FLOW-6: Handle DOCX files — extract paragraphs
    if ext == ".docx":
        if not HAS_DOCX:
            raise ValueError("DOCX support unavailable. Install python-docx.")
        document = DocxDocument(io.BytesIO(data))
        parts = [p.text.strip() for p in document.paragraphs if p.text and p.text.strip()]

        return "\n".join(parts), metadata

    # FLOW-7: Reject legacy DOC format
    if ext == ".doc":
        raise ValueError("Legacy .doc files are not directly supported. Please convert to .docx and upload again.")

    # FLOW-8: Handle PPTX files — extract text from all shapes per slide
    if ext == ".pptx":
        if not HAS_PPTX:
            raise ValueError("PPTX support unavailable. Install python-pptx.")
        presentation = Presentation(io.BytesIO(data))
        slides = []
        for idx, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    val = shape.text.strip()
                    if val:
                        texts.append(val)
            if texts:
                slides.append(f"[Slide {idx}]\n" + "\n".join(texts))

        return "\n\n".join(slides), metadata

    # FLOW-9: Reject legacy PPT format
    if ext == ".ppt":
        raise ValueError("Legacy .ppt files are not directly supported. Please convert to .pptx and upload again.")

    # FLOW-10: Handle image files — use OCR to extract text
    if ext in {".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp"}:
        if not HAS_OCR:
            raise ValueError("OCR support unavailable. Install pillow + pytesseract and configure Tesseract binary.")
        image = Image.open(io.BytesIO(data))
        text = pytesseract.image_to_string(image)

        return text.strip(), metadata

    # FLOW-11: Reject any unrecognized extension
    raise ValueError(f"Unsupported file type: {ext or 'unknown'}")
# =========== FUNCTION ===========


# =========== FUNCTION ===========
# ROLE: Split text into fixed-size chunks with overlap
def chunk_text(text: str, max_chars: int = 1500, overlap: int = 200) -> List[str]:
    ''' Return list of text chunks, respecting max_chars and overlap settings '''

    # FLOW-1: Normalize whitespace and check if empty
    cleaned = " ".join((text or "").split())
    if not cleaned:
        return []

    # FLOW-2: Return single chunk if text fits in max size
    if len(cleaned) <= max_chars:
        return [cleaned]

    # FLOW-3: Split text into chunks with sliding window
    chunks: List[str] = []
    start = 0
    while start < len(cleaned):
        end = min(start + max_chars, len(cleaned))
        chunk = cleaned[start:end].strip()

        # FLOW-4: Add chunk if not empty
        if chunk:
            chunks.append(chunk)

        # FLOW-5: Break if reached end of text
        if end >= len(cleaned):
            break

        # FLOW-6: Move start position with overlap for next window
        start = max(0, end - overlap)

    return chunks
# =========== FUNCTION ===========


# =========== FUNCTION ===========
# ROLE: Build chunk metadata records from extracted text
def build_upload_chunks(
    *,
    filename: str,
    display_title: Optional[str],
    extracted_text: str,
    url_hint: str = "",
    category: str = "uploaded",
) -> List[Dict[str, str]]:
    ''' Return list of chunk dicts ready for embedding and storage '''

    # FLOW-1: Determine display title — prefer display_title, then filename
    title = (display_title or filename or "Uploaded Knowledge").strip()

    # FLOW-2: Split extracted text into chunks
    blocks = chunk_text(extracted_text, max_chars=1500, overlap=200)

    # FLOW-3: Wrap each chunk with metadata
    chunks: List[Dict[str, str]] = []
    for idx, block in enumerate(blocks):
        chunks.append(
            {
                "title": title,
                "text": block,
                "url": url_hint,
                "category": category,
                "section_type": f"upload_chunk_{idx + 1}",
                "metadata": {
                    "filename": filename,
                    "chunk_index": idx,
                },
            }
        )

    return chunks
# =========== FUNCTION ===========


# =========== FUNCTION ===========
# ROLE: Persist extracted text to sidecar file
def save_extracted_text(upload_id, text) -> Optional[str]:
    ''' Write extracted text to disk and return path, or None if failed '''

    # FLOW-1: Ensure sidecar directory exists
    try:
        os.makedirs(_SIDECAR_DIR, exist_ok=True)

        # FLOW-2: Build sidecar filename using upload ID
        path = os.path.join(_SIDECAR_DIR, 'upload_' + str(upload_id) + '_extracted.txt')

        # FLOW-3: Write text to file
        with open(path, 'w', encoding='utf-8') as fh:
            fh.write(text)

        return path

    except Exception:
        return None
# =========== FUNCTION ===========


# =========== FUNCTION ===========
# ROLE: Load extracted text from sidecar file
def load_extracted_text(path) -> Optional[str]:
    ''' Read extracted text from file, or None if not found or read failed '''

    # FLOW-1: Return None if path missing or not a file
    if not path or not os.path.isfile(path):
        return None

    # FLOW-2: Read and return file contents
    try:
        with open(path, 'r', encoding='utf-8') as fh:
            return fh.read()
    except Exception:
        return None
# =========== FUNCTION ===========